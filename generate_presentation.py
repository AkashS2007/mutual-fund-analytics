"""
Generate 12-slide PowerPoint Presentation for Bluestock MF Capstone Project.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colors ───────────────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
BLUE    = RGBColor(0x00, 0xA3, 0xE0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF0, 0xF4, 0xF8)
DGRAY   = RGBColor(0x44, 0x44, 0x44)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
ORANGE  = RGBColor(0xF3, 0x9C, 0x12)

W = Inches(13.33)   # Widescreen width
H = Inches(7.5)     # Widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # Blank layout

# ── Helpers ───────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def kpi_card(slide, x, y, value, label, val_color=BLUE):
    box(slide, x, y, 2.8, 1.4, WHITE)
    txt(slide, value, x+0.1, y+0.1, 2.6, 0.7, size=28, bold=True,
        color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x+0.1, y+0.85, 2.6, 0.45, size=11,
        color=DGRAY, align=PP_ALIGN.CENTER)

def bullet_block(slide, items, x, y, w, size=13, color=DGRAY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 13.33, 0.08, BLUE)
box(s, 0, 7.42, 13.33, 0.08, BLUE)

txt(s, "BLUESTOCK FINTECH  |  BATCH MJ28", 0, 0.5, 13.33, 0.5,
    size=13, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, "Capstone Project I", 0, 1.3, 13.33, 0.8,
    size=22, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, "Mutual Fund Analytics Platform", 0, 2.0, 13.33, 1.2,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "End-to-End Data Analytics | ETL · SQL · EDA · ML · Power BI",
    0, 3.4, 13.33, 0.6, size=16, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

box(s, 3.5, 4.2, 6.33, 0.04, BLUE)

txt(s, "Team: akashs3882  ·  bandarivyshnavi2  ·  vhedau47  ·  Kota Naga Raviteja  ·  divyamadiraju2006",
    0, 4.5, 13.33, 0.5, size=12, color=LGRAY, align=PP_ALIGN.CENTER)
txt(s, "June 2026", 0, 5.1, 13.33, 0.4,
    size=13, color=BLUE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Objective
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "Problem & Objective", 0.4, 0.2, 10, 0.7,
    size=32, bold=True, color=WHITE)

txt(s, "PROBLEM", 0.5, 1.4, 5.5, 0.5, size=14, bold=True, color=NAVY)
bullet_block(s, [
    "India's MF industry generates massive data — NAV, AUM,",
    "  SIP flows, investor transactions — but lacks unified analytics",
    "No single platform connects fund performance, investor",
    "  behaviour, and market benchmarks",
    "AMCs struggle to identify at-risk SIP investors early",
], 0.5, 1.9, 5.8, size=13)

txt(s, "OBJECTIVE", 7.2, 1.4, 5.5, 0.5, size=14, bold=True, color=NAVY)
bullet_block(s, [
    "Build end-to-end MF analytics platform on 10 AMFI datasets",
    "Design SQLite star-schema and load 1,42,000+ records",
    "Compute financial metrics: CAGR, Sharpe, Alpha, VaR",
    "Build interactive Power BI dashboard (4 pages)",
    "Deliver actionable insights for AMCs and investors",
], 7.2, 1.9, 5.8, size=13)

box(s, 0.5, 5.6, 12.3, 1.5, LGRAY)
txt(s, "Scope: 40 Schemes  |  10 AMCs  |  64,320 NAV Records  |  32,778 Transactions  |  2022-2026",
    0.5, 5.75, 12.3, 0.5, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(s, "Rs. 81 Lakh Crore AUM  |  26.12 Crore Folios  |  Rs. 31,002 Cr SIP ATH",
    0.5, 6.2, 12.3, 0.5, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Data Sources
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "Data Sources", 0.4, 0.2, 10, 0.7, size=32, bold=True, color=WHITE)

datasets = [
    ("01", "fund_master",            "40 schemes | metadata, risk, expense ratio"),
    ("02", "nav_history",            "46,000 rows | daily NAV 2022-2026"),
    ("03", "aum_by_fund_house",      "90 rows | quarterly AUM by AMC"),
    ("04", "monthly_sip_inflows",    "48 rows | monthly SIP flows"),
    ("05", "category_inflows",       "144 rows | category-wise net inflows"),
    ("06", "industry_folio_count",   "21 rows | total industry folios"),
    ("07", "scheme_performance",     "40 rows | Sharpe, Alpha, Beta, drawdown"),
    ("08", "investor_transactions",  "32,778 rows | demographics + amounts"),
    ("09", "portfolio_holdings",     "322 rows | sector/stock weights"),
    ("10", "benchmark_indices",      "8,050 rows | Nifty 50 & Nifty 100"),
]

cols = [datasets[:5], datasets[5:]]
for ci, col in enumerate(cols):
    cx = 0.4 + ci * 6.5
    for ri, (num, name, desc) in enumerate(col):
        ry = 1.3 + ri * 1.1
        box(s, cx, ry, 6.1, 0.95, LGRAY)
        txt(s, num, cx+0.1, ry+0.05, 0.5, 0.4, size=14, bold=True, color=BLUE)
        txt(s, name, cx+0.6, ry+0.05, 5.2, 0.4, size=12, bold=True, color=NAVY)
        txt(s, desc, cx+0.6, ry+0.45, 5.2, 0.4, size=10, color=DGRAY)

txt(s, "Source: Bluestock Fintech Platform + mfapi.in Live API",
    0, 7.1, 13.33, 0.35, size=11, color=BLUE,
    align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "System Architecture & ETL Pipeline", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

stages = [
    (BLUE,   "EXTRACT",   "10 CSVs + Live\nmfapi.in API"),
    (NAVY,   "TRANSFORM", "Clean · Validate\nStandardise"),
    (GREEN,  "LOAD",      "SQLite Star\nSchema DB"),
    (ORANGE, "ANALYSE",   "EDA · Metrics\nML Models"),
    (BLUE,   "VISUALISE", "Power BI\n4-Page Dashboard"),
]
for i, (color, stage, desc) in enumerate(stages):
    x = 0.4 + i * 2.55
    box(s, x, 1.4, 2.2, 1.5, color)
    txt(s, stage, x, 1.4, 2.2, 0.7, size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, x, 2.0, 2.2, 0.9, size=11,
        color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", x+2.2, 1.85, 0.35, 0.6, size=22,
            bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txt(s, "Database Schema", 0.4, 3.2, 5, 0.5, size=14, bold=True, color=NAVY)
tables = [
    ("dim_fund",          "40 rows",    "amfi_code PK"),
    ("fact_nav",          "64,320 rows","amfi_code FK · date · nav"),
    ("fact_transactions", "32,778 rows","amfi_code FK · amount"),
    ("fact_performance",  "40 rows",    "Sharpe · Alpha · Beta"),
    ("fact_aum",          "90 rows",    "fund_house · aum_crore"),
]
for i, (name, rows, cols_txt) in enumerate(tables):
    ry = 3.8 + i * 0.65
    box(s, 0.4, ry, 2.5, 0.55, LGRAY)
    txt(s, name, 0.5, ry+0.05, 2.3, 0.45, size=11, bold=True, color=NAVY)
    txt(s, rows, 3.0, ry+0.05, 1.8, 0.45, size=11, color=BLUE)
    txt(s, cols_txt, 4.9, ry+0.05, 4.5, 0.45, size=10, color=DGRAY)

txt(s, "Tech Stack", 7.5, 3.2, 5.5, 0.5, size=14, bold=True, color=NAVY)
stack = ["Python · Pandas · NumPy · SciPy",
         "SQLite · SQLAlchemy · Plotly",
         "Seaborn · Matplotlib · Jupyter",
         "Power BI Desktop · mfapi.in API",
         "GitHub · VS Code"]
bullet_block(s, stack, 7.5, 3.7, 5.5, size=12)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA Highlights 1
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "EDA Highlights — Industry Overview", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

kpis = [
    ("Rs. 81L Cr",   "Total Industry AUM",       BLUE),
    ("Rs. 31,002 Cr","SIP All-Time High",         GREEN),
    ("26.12 Cr",     "Total Folios (Dec 2025)",   ORANGE),
    ("40 Schemes",   "Across 10 AMCs",            NAVY),
]
for i, (val, label, color) in enumerate(kpis):
    kpi_card(s, 0.4 + i * 3.2, 1.3, val, label, color)

txt(s, "Key EDA Findings", 0.4, 2.95, 12, 0.45, size=14, bold=True, color=NAVY)
findings = [
    "SBI Mutual Fund dominates AUM at Rs. 12.5L Cr — nearly 3x its closest competitor ICICI Prudential",
    "SIP inflows grew 169% from Rs. 11,517 Cr (Jan 2022) to Rs. 31,002 Cr (Dec 2025) — all-time high",
    "Flexi Cap and Small Cap categories attract consistently highest net inflows across all months",
    "25-35 age group drives 42% of all SIP transactions; male investors outnumber female 2:1",
    "Maharashtra & Karnataka lead state-wise SIP volumes; T30 cities = 65% of total investment",
    "Folios doubled from 13.26 Cr to 26.12 Cr — India's mutual fund penetration growing rapidly",
]
bullet_block(s, findings, 0.4, 3.4, 12.5, size=12)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — EDA Highlights 2
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "EDA Highlights — NAV & Correlation Analysis", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

txt(s, "NAV Trend Observations", 0.4, 1.3, 6, 0.45, size=14, bold=True, color=NAVY)
nav_points = [
    "All 40 funds showed strong NAV appreciation during 2023 bull run",
    "Small Cap funds gained 35-45% in 2023 — highest across all categories",
    "Mid-2024 saw a correction of 8-12% for equity funds, recovering by Q4",
    "Debt funds showed stable, low-volatility NAV growth throughout",
    "2025 saw renewed bull momentum across all equity categories",
]
bullet_block(s, nav_points, 0.4, 1.85, 5.8, size=12)

txt(s, "Correlation Matrix Insights", 7.0, 1.3, 6, 0.45, size=14, bold=True, color=NAVY)
corr_points = [
    "Equity funds show high positive correlation (> 0.80) with each other",
    "Large Cap funds track Nifty 50 most closely (Beta near 1.0)",
    "Small Cap funds have lower correlation with Large Cap peers",
    "Debt funds show near-zero correlation with equity — good diversifiers",
    "Gilt funds show negative correlation during equity downturns",
]
bullet_block(s, corr_points, 7.0, 1.85, 6.0, size=12)

box(s, 0.4, 5.1, 12.5, 1.9, LGRAY)
txt(s, "Sector Allocation Insight", 0.6, 5.2, 12, 0.4,
    size=13, bold=True, color=NAVY)
txt(s, "Financial Services (28%) and Information Technology (14%) dominate equity fund allocations, "
       "comprising 42% of total portfolio weight. Healthcare and Consumer Goods follow at 11% and 9% respectively. "
       "Mid Cap and Small Cap funds show higher HHI concentration scores (>0.15), indicating concentrated sector bets.",
    0.6, 5.6, 12.1, 1.3, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Performance Metrics 1
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "Performance Metrics — Risk & Return", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

metrics = [
    ("CAGR",           "Compound Annual Growth Rate\n(1yr, 3yr, 5yr)",
     "Small Cap avg 3yr CAGR: 24.3%\nLarge Cap avg 3yr CAGR: 16.1%"),
    ("Sharpe Ratio",   "Excess return per unit of\ntotal risk (Rf = 6.5%)",
     "Top fund Sharpe: 1.82\nIndustry avg: 0.94"),
    ("Sortino Ratio",  "Excess return per unit of\ndownside risk only",
     "Better than Sharpe for\nasymmetric return funds"),
    ("Alpha & Beta",   "OLS regression vs\nNifty 100 benchmark",
     "Top Alpha fund: +8.4%\nAvg Beta: 0.91"),
    ("Max Drawdown",   "Peak-to-trough decline\nfor each fund",
     "Worst drawdown: -31.2%\n(Small Cap, mid-2024)"),
    ("VaR / CVaR",     "95% Value at Risk &\nConditional VaR",
     "Avg VaR (95%): -1.84%/day\nCVaR: -2.61%/day"),
]
for i, (name, formula, result) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.3 + row * 2.8
    box(s, x, y, 4.0, 2.5, LGRAY)
    txt(s, name, x+0.15, y+0.1, 3.7, 0.5, size=14, bold=True, color=NAVY)
    txt(s, formula, x+0.15, y+0.6, 3.7, 0.9, size=11, color=DGRAY)
    box(s, x+0.1, y+1.55, 3.8, 0.8, BLUE)
    txt(s, result, x+0.2, y+1.6, 3.6, 0.7, size=10, color=WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Performance Metrics 2 (Scorecard)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "Fund Scorecard (0-100 Composite Score)", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

txt(s, "Scorecard Methodology", 0.4, 1.25, 6, 0.45, size=14, bold=True, color=NAVY)
components = [
    ("30%", "3-Year CAGR Rank",         "Higher is better"),
    ("25%", "Sharpe Ratio Rank",         "Higher is better"),
    ("20%", "Alpha Rank",                "Higher is better"),
    ("15%", "Expense Ratio Rank",        "Lower is better (inverse)"),
    ("10%", "Max Drawdown Rank",         "Lower is better (inverse)"),
]
for i, (wt, comp, note) in enumerate(components):
    ry = 1.8 + i * 0.75
    box(s, 0.4, ry, 0.7, 0.6, BLUE)
    txt(s, wt, 0.4, ry+0.1, 0.7, 0.4, size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, comp, 1.2, ry+0.05, 3.5, 0.3, size=12, bold=True, color=NAVY)
    txt(s, note, 1.2, ry+0.33, 3.5, 0.25, size=10, color=DGRAY, italic=True)

txt(s, "Top 5 Funds by Score", 7.0, 1.25, 6, 0.45, size=14, bold=True, color=NAVY)
top5 = [
    ("1st", "SBI Small Cap Fund - Direct",       "87.4"),
    ("2nd", "Nippon India Large Cap - Direct",    "82.1"),
    ("3rd", "ICICI Pru Bluechip - Direct",        "79.6"),
    ("4th", "Kotak Bluechip - Direct",            "76.3"),
    ("5th", "Axis Bluechip - Direct",             "74.8"),
]
for i, (rank, name, score) in enumerate(top5):
    ry = 1.8 + i * 0.9
    box(s, 7.0, ry, 0.7, 0.75, NAVY)
    txt(s, rank, 7.0, ry+0.15, 0.7, 0.45, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 7.75, ry, 4.3, 0.75, LGRAY)
    txt(s, name, 7.85, ry+0.05, 3.8, 0.4, size=11, color=NAVY, bold=True)
    box(s, 12.1, ry, 0.8, 0.75, GREEN)
    txt(s, score, 12.1, ry+0.15, 0.8, 0.45, size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Advanced Risk Insights", 0.4, 5.85, 12.5, 0.4, size=13, bold=True, color=NAVY)
box(s, 0.4, 6.3, 12.5, 0.9, LGRAY)
txt(s, "18% of SIP investors flagged as at-risk (avg gap > 35 days)  |  "
       "Rolling 90-day Sharpe dipped below 0 in mid-2024 for all 5 key funds  |  "
       "Small Cap VaR: -2.4%/day vs Large Cap: -1.3%/day",
    0.6, 6.38, 12.1, 0.75, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Dashboard Screenshot 1
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "Power BI Dashboard — Page 1 & 2", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

box(s, 0.3, 1.2, 6.0, 4.5, LGRAY)
txt(s, "Page 1: Industry Overview", 0.5, 1.3, 5.6, 0.45,
    size=13, bold=True, color=NAVY)
txt(s, "KPI Cards: Total AUM · SIP Inflows · Folios · Schemes\n"
       "Line Chart: Industry AUM Trend 2022-2025\n"
       "Bar Chart: AUM by AMC (SBI dominance highlighted)\n"
       "Slicers: Year · Fund House",
    0.5, 1.85, 5.6, 2.0, size=11, color=DGRAY)
box(s, 0.5, 4.0, 5.6, 1.5, RGBColor(0xE8, 0xF4, 0xFD))
txt(s, "Rs. 81L Cr Total AUM | Rs. 31K Cr SIP ATH\n26.12 Cr Folios | 1,908 Schemes",
    0.6, 4.1, 5.4, 1.3, size=12, color=NAVY, align=PP_ALIGN.CENTER)

box(s, 7.0, 1.2, 6.0, 4.5, LGRAY)
txt(s, "Page 2: Fund Performance", 7.2, 1.3, 5.6, 0.45,
    size=13, bold=True, color=NAVY)
txt(s, "Scatter Plot: Return vs Risk (bubble = AUM)\n"
       "Sortable Scorecard Table (all 40 funds)\n"
       "NAV Line vs Nifty 50 / Nifty 100 benchmark\n"
       "Slicers: Fund House · Category · Plan\n"
       "Drill-through: Fund → NAV Detail Page",
    7.2, 1.85, 5.6, 2.2, size=11, color=DGRAY)
box(s, 7.2, 4.0, 5.6, 1.5, RGBColor(0xE8, 0xF4, 0xFD))
txt(s, "Top Sharpe: 1.82 | Top Alpha: +8.4%\nBenchmark Tracking Error: 4.2%",
    7.3, 4.1, 5.4, 1.3, size=12, color=NAVY, align=PP_ALIGN.CENTER)

txt(s, "Bluestock Colour Theme Applied | Tooltips on all charts | Drill-through enabled",
    0, 6.9, 13.33, 0.45, size=11, color=BLUE,
    align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Dashboard Screenshot 2
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "Power BI Dashboard — Page 3 & 4", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

box(s, 0.3, 1.2, 6.0, 4.5, LGRAY)
txt(s, "Page 3: Investor Analytics", 0.5, 1.3, 5.6, 0.45,
    size=13, bold=True, color=NAVY)
txt(s, "Bar Chart: Transaction amount by state (Top 15)\n"
       "Donut: SIP / Lumpsum / Redemption split\n"
       "Bar: Age group vs avg SIP amount\n"
       "Line: Monthly transaction volume trend\n"
       "Slicers: State · Age Group · City Tier",
    0.5, 1.85, 5.6, 2.2, size=11, color=DGRAY)
box(s, 0.5, 4.0, 5.6, 1.5, RGBColor(0xE8, 0xF4, 0xFD))
txt(s, "SIP: 68% | Lumpsum: 22% | Redemption: 10%\nTop State: Maharashtra",
    0.6, 4.1, 5.4, 1.3, size=12, color=NAVY, align=PP_ALIGN.CENTER)

box(s, 7.0, 1.2, 6.0, 4.5, LGRAY)
txt(s, "Page 4: SIP & Market Trends", 7.2, 1.3, 5.6, 0.45,
    size=13, bold=True, color=NAVY)
txt(s, "Dual-axis: SIP Inflow (bar) + Nifty 50 (line)\n"
       "Category Inflow Heatmap (months vs categories)\n"
       "Top 5 Categories by Net Inflow FY25\n"
       "SIP YoY Growth Trend Annotation\n"
       "Market correlation indicators",
    7.2, 1.85, 5.6, 2.2, size=11, color=DGRAY)
box(s, 7.2, 4.0, 5.6, 1.5, RGBColor(0xE8, 0xF4, 0xFD))
txt(s, "SIP-Market Correlation: 0.74\nFlexi Cap: Highest FY25 Inflow Category",
    7.3, 4.1, 5.4, 1.3, size=12, color=NAVY, align=PP_ALIGN.CENTER)

txt(s, "Exported as .pbix | PDF Export | 4 PNG Screenshots for Final Report",
    0, 6.9, 13.33, 0.45, size=11, color=BLUE,
    align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Findings
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
box(s, 0, 0, 13.33, 1.1, NAVY)
txt(s, "Key Findings & Recommendations", 0.4, 0.2, 12, 0.7,
    size=32, bold=True, color=WHITE)

findings_data = [
    (BLUE,   "Industry Growth",
     "AUM doubled to Rs. 81L Cr. SIP inflows hit Rs. 31,002 Cr ATH. "
     "Folios doubled to 26.12 Cr in just 4 years."),
    (GREEN,  "Top Performers",
     "Small Cap funds outperform Large Cap by 8-12% over 3 years. "
     "Direct plans consistently beat Regular plans by 1-1.5% annually."),
    (ORANGE, "Investor Behaviour",
     "18% of SIP investors are at-risk. 25-35 age group drives 42% of SIPs. "
     "B30 cities represent 35% — an underserved growth opportunity."),
    (NAVY,   "Risk Insights",
     "Small Cap VaR is 2x Large Cap. Mid-2024 saw max drawdown of -31.2%. "
     "Rolling Sharpe recovered strongly by Q1 2025."),
    (RGBColor(0x8E, 0x44, 0xAD), "Recommendations",
     "Target at-risk SIP investors with automated reminders. "
     "Focus digital campaigns on B30 cities. Disclose HHI concentration risk."),
]
for i, (color, title_t, desc) in enumerate(findings_data):
    x = 0.4 + (i % 3) * 4.3
    y = 1.35 + (i // 3) * 2.8
    box(s, x, y, 0.12, 2.4, color)
    box(s, x+0.18, y, 3.9, 2.4, LGRAY)
    txt(s, title_t, x+0.28, y+0.1, 3.7, 0.5, size=13, bold=True, color=color)
    txt(s, desc,    x+0.28, y+0.6, 3.7, 1.6, size=11, color=DGRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
box(s, 0, 0, 13.33, 0.08, BLUE)
box(s, 0, 7.42, 13.33, 0.08, BLUE)

txt(s, "Thank You", 0, 1.5, 13.33, 1.5,
    size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Capstone Project I — Mutual Fund Analytics Platform",
    0, 3.0, 13.33, 0.6, size=18, color=BLUE,
    align=PP_ALIGN.CENTER, italic=True)

box(s, 3.0, 3.8, 7.33, 0.04, BLUE)

deliverables = [
    "data_ingestion.py  |  live_nav_fetch.py  |  data_cleaning.py",
    "bluestock_mf.db  |  schema.sql  |  queries.sql",
    "EDA_Analysis.ipynb  |  Performance_Analytics.ipynb  |  Advanced_Analytics.ipynb",
    "bluestock_mf_dashboard.pbix  |  Final_Report.pdf",
    "README.md  |  run_pipeline.py  |  GitHub v1.0 Tag",
]
for i, d in enumerate(deliverables):
    txt(s, d, 0, 4.1 + i*0.5, 13.33, 0.45,
        size=11, color=LGRAY, align=PP_ALIGN.CENTER)

txt(s, "Team MJ28  |  Bluestock Fintech Internship  |  June 2026",
    0, 6.8, 13.33, 0.45, size=12, color=BLUE, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────
os.makedirs("reports", exist_ok=True)
prs.save("reports/Bluestock_MF_Presentation.pptx")
print("✅ Presentation saved → reports/Bluestock_MF_Presentation.pptx")